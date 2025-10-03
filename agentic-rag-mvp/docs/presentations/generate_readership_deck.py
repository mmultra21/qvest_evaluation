from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os

SLIDES = [
    {
        "title": "Increase Student Readership — Overview",
        "bullets": [
            "Use case: increase book readership among students through librarian recommendations and student visibility",
            "Goal: boost engagement, measure impact, and identify top readers"
        ],
        "notes": "High-level summary of campaign goals and expected outcomes."
    },
    {
        "title": "Proposed Solution",
        "bullets": [
            "Librarian curates recommendations using historical data and an evaluation tool",
            "Students see personalized recommendations and can opt in to campaigns",
            "Lightweight web UI + backend recommendation engine; optional Gradio proof-of-concept"
        ],
        "notes": "Describe architecture: data -> model/rules -> UI -> feedback loop."
    },
    {
        "title": "Workflow (Librarian & Student)",
        "bullets": [
            "1) Librarian uploads or selects historical circulation & interaction data",
            "2) System analyzes historical signals (ratings, checkouts, reading time) to surface candidates",
            "3) Librarian reviews, evaluates, and tags recommendations",
            "4) Student receives visible librarian recommendations in UI; can accept/ignore or rate",
            "5) Feedback is stored for continuous improvement"
        ],
        "notes": "Assumes librarian has access to historical data and evaluation tools."
    },
    {
        "title": "Recommendation Criteria",
        "bullets": [
            "Historical popularity among similar students (cohort affinity)",
            "Relevance to student interests and reading level",
            "Freshness (new releases) vs. long-tail gems balance",
            "Diversity and inclusive representation",
            "Librarian's expert tags and manual overrides"
        ],
        "notes": "Criteria mix of automated signals and librarian judgement."
    },
    {
        "title": "Campaign Mechanics & Incentives",
        "bullets": [
            "Weekly reading challenge with a prize for 'Top Reader of the Week'",
            "Track completed readings, time spent, reviews/ratings",
            "Badging and progress indicators visible in student profile",
            "Email/dashboards for librarian showing campaign progress"
        ],
        "notes": "Incentives encourage participation; track progress centrally."
    },
    {
        "title": "Out-of-Scope for MVP (can be added later)",
        "bullets": [
            "Advanced metric generation and long-term student success modeling",
            "Full analytics dashboards and BI integrations",
            "Complex A/B experimentation infrastructure",
            "Automated content generation beyond recommendations"
        ],
        "notes": "Keep MVP focused; mark advanced features for phase 2."
    },
    {
        "title": "Monitoring & Metrics",
        "bullets": [
            "Adoption: % of students who view and interact with recommendations",
            "Engagement: books read, time spent, ratings submitted",
            "Campaign success: delta in readership vs baseline",
            "Top reader leaderboard and weekly prize delivery metrics"
        ],
        "notes": "Define baseline and success thresholds before campaign start."
    },
    {
        "title": "Scaling Plan (if approved)",
        "bullets": [
            "Phase 1 (MVP): single-region deployment, minimal infra, run on a single VM or managed instance",
            "Phase 2: containerize (Docker), add CI/CD, autoscaling backend on AWS ECS/EKS or GCP GKE",
            "Phase 3: multi-region, production-grade monitoring, managed DB (RDS/Cloud SQL), CDN for static assets"
        ],
        "notes": "Progressive approach reduces risk and cost."
    },
    {
        "title": "Estimated Timeline & Tasks",
        "bullets": [
            "Weeks 0-2: Requirements, data access, prototype UI (Gradio or simple Flask/Streamlit)",
            "Weeks 3-5: Build recommendation engine and librarian evaluation UI",
            "Weeks 6-8: Integrate student UI, telemetry, and deployment scripts",
            "Weeks 9-12: Pilot, iterate on feedback, baseline comparisons"
        ],
        "notes": "Timeline assumes small team (1-2 devs + librarian + product owner)."
    },
    {
        "title": "Hypothetical Costs (very rough)",
        "bullets": [
            "Initial MVP (dev work): 2-3 person-months (~$30k-$60k depending on rates)",
            "Cloud infra (small): $50-$300/month for dev/test; $300-$2k+/month for production depending on scale",
            "Progressive growth: add managed DB, load balancer, and autoscaling as needed"
        ],
        "notes": "Provide a conservative cost band; refine after architecture review."
    },
    {
        "title": "Rollout Strategy",
        "bullets": [
            "Pilot with a single librarian cohort and small student group",
            "Measure metrics, collect qualitative feedback, iterate",
            "Gradual expansion by school/grade with repeated measurement",
            "Full production rollout after meeting success thresholds"
        ],
        "notes": "Pilot-first reduces risk."
    },
    {
        "title": "Call to Action",
        "bullets": [
            "Approve pilot: provide access to historical circulation data and a librarian partner",
            "Allocate 4-6 weeks of development resources to deliver MVP",
            "Define success metrics and prize criteria for top reader incentive"
        ],
        "notes": "Clear asks to stakeholders to get started."
    }
]


def make_deck(out_path):
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Readership Campaign: Increase Student Reading"
    subtitle.text = "Use case, workflow, criteria, incentives, and scaling plan"

    # add content slides
    for s in SLIDES:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.shapes.placeholders[1].text_frame
        title.text = s['title']
        body.text = s['bullets'][0]
        for b in s['bullets'][1:]:
            p = body.add_paragraph()
            p.text = b
            p.level = 1
        # add notes
        if 'notes' in s and s['notes']:
            slide.notes_slide.notes_text_frame.text = s['notes']

    # Add a process flow image slide (if a PNG is available). We ship an SVG; convert to PNG first
    flow_png = os.path.join(os.path.dirname(out_path), 'process_flow.png')
    flow_svg = os.path.join(os.path.dirname(out_path), 'process_flow.svg')
    if os.path.exists(flow_png):
        img_slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
        slide = prs.slides.add_slide(img_slide_layout)
        title = slide.shapes.title if slide.shapes.title else None
        if title:
            title.text = 'Application Process Flow'
        # add picture full-bleed-ish
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        slide.shapes.add_picture(flow_png, left, top, width=width)
    else:
        # If no PNG, add a slide that points to the SVG and instructions
        info_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(info_layout)
        slide.shapes.title.text = 'Application Process Flow (image not embedded)'
        body = slide.shapes.placeholders[1].text_frame
        if os.path.exists(flow_svg):
            body.text = 'The process flow diagram is available as SVG at:`{}`.\nConvert it to PNG and rerun the generator to embed it.'.format(flow_svg)
        else:
            body.text = 'Process flow SVG not found. Please create process_flow.png in the presentations folder to embed the image.'

    # final save
    out_dir = os.path.dirname(out_path)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir)
    prs.save(out_path)
    print('Deck created at', out_path)


if __name__ == '__main__':
    out = os.path.join('agentic-rag-mvp','docs','presentations','readership_campaign_deck.pptx')
    make_deck(out)
